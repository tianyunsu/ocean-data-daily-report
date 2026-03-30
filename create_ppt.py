# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Colors
RED = RGBColor(200, 16, 46)
GOLD = RGBColor(255, 215, 0)
DARK_RED = RGBColor(139, 0, 0)
LIGHT_BEIGE = RGBColor(255, 245, 230)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(51, 51, 51)

def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, color=RED, size=44):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = text
    title_frame.paragraphs[0].font.size = Pt(size)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = color
    title_frame.paragraphs[0].font.name = 'SimHei'
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_line(slide):
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(2.5), Pt(4))
    line.line.color.rgb = GOLD
    line.line.width = Pt(4)

def add_text(slide, text, left, top, width, height, color=BLACK, size=20, bold=False):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(size)
    text_frame.paragraphs[0].font.bold = bold
    text_frame.paragraphs[0].font.color.rgb = color
    text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
    text_frame.word_wrap = True

def add_bullet(slide, text, left, top, width, height, color=BLACK, size=20):
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = text
    p.level = 0
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    text_frame.word_wrap = True

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide1, RED)
add_text(slide1, "树立和践行正确政绩观", 0.5, 1.5, 9, 1, WHITE, 60, True)
add_text(slide1, "学习教育工作部署", 0.5, 2.5, 9, 1, WHITE, 56, True)
add_text(slide1, "2026年3月", 3.5, 4.5, 3, 0.5, LIGHT_BEIGE, 24, False)

# Slide 2: Background
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide2, LIGHT_BEIGE)
add_title(slide2, "一、会议背景", RED, 44)
add_line(slide2)
content = """自然资源部召开党组会议，深入学习贯彻习近平总书记关于树立和践行正确政绩观的重要论述和重要指示精神，传达学习中办印发的《关于在全党开展树立和践行正确政绩观学习教育的通知》和中央党的建设工作领导小组会议精神。

自然资源部第一海洋研究所召开党委（扩大）会议，传达学习中央通知及自然资源部党组要求，动员部署全所开展树立和践行正确政绩观学习教育工作。"""
add_text(slide2, content, 0.8, 1.8, 8.4, 3.5, BLACK, 18, False)

# Slide 3: Significance
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide3, LIGHT_BEIGE)
add_title(slide3, "二、重要意义", RED, 44)
add_line(slide3)
y_pos = 1.8
bullets = [
    "贯彻落实党的二十届四中全会战略部署的重要举措",
    "践行党的根本宗旨、夯实党的执政根基的必然要求",
    "推进党和国家事业发展、全面从严治党的关键环节",
    "2026年党建工作的重要任务"
]
for bullet in bullets:
    add_bullet(slide3, bullet, 0.8, y_pos, 8.4, 0.8, BLACK, 22)
    y_pos += 0.9

# Slide 4: Requirements
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide4, LIGHT_BEIGE)
add_title(slide4, "三、总要求", RED, 44)
add_line(slide4)
add_text(slide4, "立党为公、为民造福、科学决策、真抓实干", 1.5, 2, 7, 1, RED, 32, True)
y_pos = 3.2
sub_bullets = [
    "坚持学查改一体推进",
    "结合'学论述谋发展见行动'活动",
    "确保学习教育取得实效",
    "将正确政绩观融入科研管理和科技创新全过程"
]
for bullet in sub_bullets:
    add_bullet(slide4, bullet, 1.5, y_pos, 7, 0.6, BLACK, 20)
    y_pos += 0.7

# Slide 5: Tasks
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide5, LIGHT_BEIGE)
add_title(slide5, "四、重点任务", RED, 44)
add_line(slide5)
y_pos = 1.8
tasks = [
    "深学细悟：把牢正确方向，深入学习领会习近平总书记重要论述精神",
    "联系实际：深挖细查问题，结合实际工作查找问题",
    "从严从实：推动整改整治，力戒形式主义"
]
for task in tasks:
    add_bullet(slide5, task, 0.8, y_pos, 8.4, 1.1, BLACK, 20)
    y_pos += 1.3

# Slide 6: Rectification
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide6, LIGHT_BEIGE)
add_title(slide6, "五、整治重点", RED, 44)
add_line(slide6)
y_pos = 1.8
rectifications = [
    "重点整治规划编制中的形式主义问题",
    "化解自然资源领域信访问题",
    "排查整改影响科研创新发展的突出问题",
    "解决群众反映强烈的急难愁盼问题"
]
for item in rectifications:
    add_bullet(slide6, item, 0.8, y_pos, 8.4, 0.9, BLACK, 20)
    y_pos += 1.0

# Slide 7: Work Requirements
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide7, LIGHT_BEIGE)
add_title(slide7, "六、工作要求", RED, 44)
add_line(slide7)
y_pos = 1.8
requirements = [
    "提高政治站位，深刻认识学习教育的重大意义",
    "精心组织实施，各级党组织负起政治责任",
    "坚持开门教育，广泛听取群众意见建议",
    "力戒形式主义，确保学习教育不偏不空、不走过场",
    "以正确政绩观引领自然资源工作提质量上水平",
    "以清廉务实作风推动'十五五'科研工作开好局"
]
for req in requirements:
    add_bullet(slide7, req, 0.8, y_pos, 8.4, 0.75, BLACK, 18)
    y_pos += 0.85

# Slide 8: End
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_background(slide8, RED)
add_text(slide8, "谢谢观看", 2.5, 2.5, 5, 1, WHITE, 72, True)
add_text(slide8, "树立和践行正确政绩观学习教育", 2, 3.8, 6, 0.6, LIGHT_BEIGE, 28, False)

# Save
prs.save(r'c:\Users\Administrator\WorkBuddy\Claw\树立和践行正确政绩观学习教育.pptx')
print("PPT created successfully!")
