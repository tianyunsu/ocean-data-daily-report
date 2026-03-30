import sys
sys.path.insert(0, r'C:\Program Files\Python311\Lib\site-packages')
try:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(1, 1, 8, 1)
    title.text_frame.text = "Test"
    prs.save(r'c:\Users\Administrator\WorkBuddy\Claw\test.pptx')
    print("Success!")
except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
